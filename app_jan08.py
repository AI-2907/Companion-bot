
# =================================code without Next quetion button========================
import streamlit as st
import re
from langchain.chat_models import ChatOpenAI
from langchain.document_loaders import DirectoryLoader, PyPDFLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.llms import Ollama
from langchain.vectorstores import FAISS
from langchain.prompts import ChatPromptTemplate, MessagesPlaceholder
from langchain_core.runnables.history import RunnableWithMessageHistory
from langchain_community.chat_message_histories import ChatMessageHistory
from langchain.chains import create_retrieval_chain
from langchain.chains.combine_documents import create_stuff_documents_chain
from sentence_transformers import SentenceTransformer, util
from langchain_core.chat_history import BaseChatMessageHistory
from langchain.chains import create_history_aware_retriever
from PyPDF2 import PdfReader
from langchain.embeddings import HuggingFaceEmbeddings
from langchain_openai import AzureOpenAIEmbeddings
from datetime import datetime

from transformers import pipeline
import config
import langsmith
from PyPDF2 import PdfReader
import os 
from langchain_openai import AzureChatOpenAI
from langchain.chat_models import ChatOllama
from pptx import Presentation
from langchain.chains import ConversationalRetrievalChain
# from langchain.chat_models import ChatOpenAI
from langchain.vectorstores import faiss
from langchain.embeddings.openai import OpenAIEmbeddings
from langchain.chains.conversation.memory import ConversationBufferWindowMemory
from dotenv import load_dotenv
import os

# Load environment variables from .env file
load_dotenv()


os.environ["AZURE_OPENAI_API_KEY"] =os.getenv("AZURE_OPENAI_API_KEY")
os.environ["AZURE_OPENAI_ENDPOINT"] = os.getenv("AZURE_OPENAI_ENDPOINT")
os.environ["AZURE_OPENAI_API_VERSION"] = "2024-02-01"
os.environ["AZURE_OPENAI_CHAT_DEPLOYMENT_NAME"] = "gpt-35-turbo"

# Templates for user and bot messages
bot_template = '''
<div style="display: flex; align-items: center; margin-bottom: 10px;">
    <div style="flex-shrink: 0; margin-right: 10px;">
        <img src="https://uxwing.com/wp-content/themes/uxwing/download/communication-chat-call/answer-icon.png" 
             style="max-height: 50px; max-width: 50px; border-radius: 50%; object-fit: cover;">
    </div>
    <div style="background-color: #dff9d8;color: black; padding: 10px; border-radius: 10px; max-width: 75%; word-wrap: break-word; overflow-wrap: break-word; border: 1px solid #78e08f;">
        {msg}
    </div>
</div>
'''
# bot_template = '''
# <div style="display: flex; align-items: center; margin-bottom: 10px;">
#     <div style="flex-shrink: 0; margin-right: 10px;">
#         <img src="https://uxwing.com/wp-content/themes/uxwing/download/communication-chat-call/answer-icon.png" 
#              style="max-height: 50px; max-width: 50px; border-radius: 50%; object-fit: cover;">
#     </div>
#     <div style="background-color: #f1f1f1; padding: 10px; border-radius: 10px; max-width: 75%; word-wrap: break-word; overflow-wrap: break-word;">
#         {msg}
#     </div>
# </div>
# '''

user_template = '''
<div style="display: flex; align-items: center; margin-bottom: 10px; justify-content: flex-end;">
    <div style="flex-shrink: 0; margin-left: 10px;">
        <img src="https://cdn.iconscout.com/icon/free/png-512/free-q-characters-character-alphabet-letter-36051.png?f=webp&w=512" 
             style="max-height: 50px; max-width: 50px; border-radius: 50%; object-fit: cover;">
    </div>    
    <div style="background-color: #66bfff; color: black; padding: 10px; border-radius: 10px; max-width: 75%; word-wrap: break-word; overflow-wrap: break-word;">
        {msg}
    </div>
</div>
'''
# def run_app():
#     llm = initialize_llm(OPENAI_API_KEY)

# Function to prepare and split documents
# @st.cache_resource
@st.cache_data
def prepare_and_split_docs(pdf_directory):
    loader = DirectoryLoader(pdf_directory, glob="**/*.pdf", loader_cls=PyPDFLoader)
    documents = loader.load()
    splitter = RecursiveCharacterTextSplitter.from_tiktoken_encoder(
        chunk_size=512,
        chunk_overlap=256,
        disallowed_special=(),
        separators=["\n\n", "\n", " "]
    )
    split_docs = splitter.split_documents(documents)
    return split_docs

# Function to ingest documents into vector store

# def ingest_into_vectordb(split_docs):
#     # Use OpenAI embeddings with `text-embedding-ada-002`
#     embeddings=AzureOpenAIEmbeddings(model="text-embedding-3-large")
#     # embeddings = OpenAIEmbeddings(model="text-embedding-ada-002", api_key="AZURE_OPENAI_API_KEY")
    
#     # Create a FAISS vectorstore from the documents
#     db = FAISS.from_documents(split_docs, embeddings)
    
#     # Save the vectorstore locally
#     DB_FAISS_PATH = 'vectorstore/db_faiss'
#     db.save_local(DB_FAISS_PATH)
    
#     print("Documents are inserted into FAISS vectorstore")
#     return db

def ingest_into_vectordb(split_docs):
    embedding_model1 = HuggingFaceEmbeddings(model_name='sentence-transformers/all-MiniLM-L6-v2')
    Newdb = FAISS.from_documents(split_docs, embedding_model1)
    return Newdb


def get_conversation_chain(retriever,top_k=5):
    # Initialize OpenAI's ChatGPT model
    llm = AzureChatOpenAI(
        azure_deployment=os.environ["AZURE_OPENAI_CHAT_DEPLOYMENT_NAME"],  
        openai_api_version=os.environ["AZURE_OPENAI_API_VERSION"],
        temperature=0.9,
        # messages=[{"role": "user", "content": prompt}],
        max_tokens=150,
        n=1,
        stop=None,
        streaming=True 
        
     )
    contextualize_q_system_prompt = (
        "Given the chat history and the latest user question, "
        "provide a response based on the documents.Do not provide responses from outside knowledge or web sources other knowledge domains or AI use is to be politely denied If no answer is found, "
        "respond: 'I'm sorry, but I couldn’t find an answer. Could you rephrase or provide more details?'"
    )
    contextualize_q_prompt = ChatPromptTemplate.from_messages(
        [
            ("system", contextualize_q_system_prompt),
            MessagesPlaceholder("chat_history"),
            ("human", "{input}"),
        ]
    )
    history_aware_retriever = create_history_aware_retriever(
        llm, retriever, contextualize_q_prompt
    )
    system_prompt = (
        "As a chat assistant, provide accurate and relevant information  based on the provided document in 2-3 sentences.accoring to given query find matching wordsfrom query in provided documents and give response using that "
        "Answer should be correct to the point short and brief for given quetion . If no relevant information is found, respond with: "
        "'I'm sorry, but I couldn’t find an answer to your question. Could you rephrase or provide more details?other knowledge domains or AI use is to be politely denied .Do not provide responses from outside knowledge or web sources' "
        "{context}"
    )
    qa_prompt = ChatPromptTemplate.from_messages(
        [
            ("system", system_prompt),
            MessagesPlaceholder("chat_history"),
            ("human", "{input}"),
        ]
    )
    question_answer_chain = create_stuff_documents_chain(llm, qa_prompt)
    rag_chain = create_retrieval_chain(history_aware_retriever, question_answer_chain)

    history_aware_retriever = create_history_aware_retriever(
        llm, retriever, contextualize_q_prompt
    )
    qa_prompt = ChatPromptTemplate.from_messages(
        [
            ("system", "Answer concisely based on the documents. {context}"),
            MessagesPlaceholder("chat_history"),
            ("human", "{input}"),
        ]
    )
    rag_chain = create_retrieval_chain(
        history_aware_retriever,
        question_answer_chain,
        # top_k=top_k,  # Limit to top-k relevant documents
        # similarity_threshold=similarity_threshold  # Set a similarity threshold
    )

    store = {}

    
#     return conversational_rag_chain
    def get_session_history(session_id: str) -> BaseChatMessageHistory:
        if session_id not in store:
            store[session_id] = ChatMessageHistory()
        return store[session_id]

    conversational_rag_chain = RunnableWithMessageHistory(
        rag_chain,
        get_session_history,
        input_messages_key="input",
        history_messages_key="chat_history",
        output_messages_key="answer",
    )
    return conversational_rag_chain

def get_relevant_context(query, retriever):
    docs = retriever.get_relevant_documents(query)
    return [doc[:500] for doc in docs]  # Truncate to 500 characters


pdf_directory="Data/T&LDocs"
split_docs = prepare_and_split_docs(pdf_directory)
vector_db = ingest_into_vectordb(split_docs)
retriever = vector_db.as_retriever()

# def initialize_conversation_chain(llm):
#     return ConversationalRetrievalChain(
#         llm=llm,
#         memory=ConversationBufferWindowMemory(k=5)  # Stores the last 5 interactions.
#     )


# Main Streamlit app
# col1, col2 = st.columns([1, 1])
# logopath = r"C:\Users\Akshada.P\OneDrive - SAKSOFT LIMITED\Documents\SAKSOFT Internal BOT\saksoftimg.png

# Set the page configuration
# st.set_page_config(page_title="Sales Companion Bot", layout="wide")

# Apply custom CSS for better styling
st.markdown("""
<style>
    body {
        background-color: #f0f4f8;  /* Light background color for the whole app */
    }
    .chat-container {
        border-radius: 10px;
        padding: 20px;
        background-color: #ADD8E6;  /* White background for chat */
        box-shadow: 0 4px 10px rgba(0, 0, 0, 0.1);
        max-width: 600px;
        margin: auto;
    }
    .user-message {
        background-color: #90EE90;  /* Light green for user messages */
        padding: 10px;
        border-radius: 10px;
        margin: 5px 0;
        text-align: left;
        border: 1px solid #78e08f;
    }
    .bot-message {
        background-color: #dff9d8;  /* Light orange for bot messages */
        padding: 10px;
        border-radius: 10px;
        margin: 5px 0;
        text-align: left;
        border: 1px solid #78e08f;
    }
    .sidebar .sidebar-content {
        background-color: #e0f7fa;  /* Light cyan for sidebar */
        padding: 10px;
        border-radius: 10px;
    }
</style>
""", unsafe_allow_html=True)
# Function to clean up page content
def clean_page_content(content):
    return ' '.join(content.split()).replace('\n', ' ')
# Initialize chat history if it doesn't exist
# Function to summarize the response
# def summarize_text(text):
#     # Placeholder for a summarization logic or API call
#     # You can replace this with your actual summarization code
#     return text[:150] + "..." if len(text) > 150 else text

if 'chat_history' not in st.session_state:
    st.session_state.chat_history = []
 # Initialize conversation chain
conversational_chain = get_conversation_chain(retriever)
st.session_state.conversational_chain = conversational_chain
# Set up the app title and welcome message
logopath = r"C:\Users\Akshada.P\OneDrive - SAKSOFT LIMITED\Documents\SAKSOFT Internal BOT\saksoftimg.png"
st.image(logopath, width=300)
# botlogo=r"C:\Users\Akshada.P\OneDrive - SAKSOFT LIMITED\Documents\SAKSOFT Internal BOT\botimagelogo..webp"
# st.image(botlogo, width=40)  # Adjust the width as needed

# col1, col2 = st.columns([1, 5])
# # with col1:
# #     st.title("Sales Companion Bot")
# with col1:
#     botlogo=r"C:\Users\Akshada.P\OneDrive - SAKSOFT LIMITED\Documents\SAKSOFT Internal BOT\botimagelogo..webp"
#     st.image(botlogo, width=40)  # Adjust the width as needed

# with col2:
st.title("🤖 Sales Companion Bot")
    # st.title("Sales Companion Bot")
# st.title("Sales Companion Bot")
# st.write("Welcome! How can I help you today?")

# Description Box
st.markdown("""
<div style="background-color: #f0f2f5; padding: 10px; border-radius: 5px;">
<h3 style="color: #333;">Welcome to the Sales Companion ChatBot!</h3>
<p style="color: #555;">This chatbot is designed to assist you with sales-related queries, provide insights, and help you prepare presentations. Just ask your question, and let me assist you!</p>
</div>
""", unsafe_allow_html=True)
st.write("Welcome! How can I help you today?")
#  Initialize session state
if "chat_history" not in st.session_state:
    st.session_state["chat_history"] = []

if "uploaded_files" not in st.session_state:
    st.session_state["uploaded_files"] = []

if "recommended_prompts" not in st.session_state:
    st.session_state["recommended_prompts"] = [
        "Give me HCL commerce case studies",
        "Summarize case studies for Logistics",
        "Create a pitch deck for a potential client in the Tranportation and Logistics sector.",
        "Tell me top3 case studies Transport and Logistics",
        "give me list of services saksoft offers"
    ]

if "active_section" not in st.session_state:
    st.session_state["active_section"] = "Chat"
if 'active_section' not in st.session_state:
    st.session_state.active_section = "Chat"  # Initial section
if 'uploaded_files' not in st.session_state:
    st.session_state.uploaded_files = []  # To store uploaded file names
# Ensure the folder exists
os.makedirs(pdf_directory, exist_ok=True)


# Sidebar with navigation menus
with st.sidebar:
    # st.title("Chatbot Navigation")
    st.sidebar.title("Assistant Portal")
    
    # Navigation buttons
    if st.button("📞 Contact Support"):
        st.session_state["active_section"] = "Contact Support"
    # elif st.session_state["active_section"] == "Contact Support":
    #  st.subheader("📞 Contact Support")
        st.markdown("**Email:** support@saksoft.com")
        st.markdown("**Phone:** +1-800-555-0199")
    if st.button("🔍 Explore More"):
        st.session_state["active_section"] = "Explore More"
    #     elif st.session_state["active_section"] == "Explore More":
    # st.subheader("🔍 Explore More")
        st.markdown("""
    - Learn more about the corporate policies.
    - Get insights on market trends.
    - Explore our latest whitepapers.
    """)
    # Button to switch to Upload Files section
    if st.sidebar.button("📂 Upload Files"):
        st.session_state.active_section = "Upload Files"

# Show Upload Files section only when active
    if st.session_state.active_section == "Upload Files":
        st.subheader("📂 Upload New Files")
    uploaded_files = st.file_uploader("Upload your files (PDF, TXT, DOCX)", accept_multiple_files=True)
    if uploaded_files:
        for uploaded_file in uploaded_files:
            # Save the uploaded file to the specified folder
            file_path = os.path.join(pdf_directory, uploaded_file.name)
            with open(file_path, "wb") as f:
                f.write(uploaded_file.getbuffer())
            st.session_state.uploaded_files.append(uploaded_file.name)
            st.success(f"File '{uploaded_file.name}' uploaded successfully! You can now ask questions based on this file.")


        # -----------------------working==================
    # if st.button("📂 Upload Files"):
    #     st.session_state["active_section"] = "Upload Files"
    # #     elif st.session_state["active_section"] == "Upload Files":
    # # st.subheader("📂 Upload New Files")
    #     uploaded_files = st.file_uploader("Upload your files (PDF, TXT, DOCX)", accept_multiple_files=True)
    #     if uploaded_files:
    #         for uploaded_file in uploaded_files:
    #             st.session_state["uploaded_files"].append(uploaded_file)
    #             st.success("Files uploaded successfully! You can now ask questions based on these files.")
    
    
    # if st.button("🕒 Chat History"):
    #     st.session_state["active_section"] = "Chat History"
    # #     elif st.session_state["active_section"] == "Chat History":
    # # st.subheader("🕒 Previous Chat History")
    #     if st.session_state["chat_history"]:
    #         for i, chat in enumerate(st.session_state["chat_history"][-5:], 1):  # Show the last 5 chats
    #             st.markdown(f"**{i}. You:** {chat['user']}")
    #             st.markdown(f"**Bot:** {chat['bot']}")
    #     else:
    #         st.markdown("No previous chat history available.")
     # Button for summarizing the response
    if st.button("💡 Recommended Prompts"):
        st.session_state["active_section"] = "Recommended Prompts"

    if st.session_state["active_section"] == "Recommended Prompts":
        # st.subheader("💡 Recommended Prompts")
        for prompt in st.session_state["recommended_prompts"]:
            if st.button(prompt):
            # Display selected prompt
                # st.markdown(f"**You selected prompt:** {prompt}")
                session_id = "user412"
            # Assuming `source_document` contains information for generating a response
                # source_document = st.session_state.get("source_document", "")
                # source_document=st.session_state.get("context_docs", [])
               # Prepare source documents if available
                # source_document = ""
                # if st.session_state.get("context_docs"):
                #     with st.expander("Source Documents"):
                #         for doc in st.session_state["context_docs"]:
                #             st.write(f"Source: {doc.metadata.get('source', 'N/A')}")
                #             source_document += doc.page_content + "\n"  # Collecting all document content
                # st.session_state.chat_history.append({"user": prompt, "bot": "", "context_docs": []})
                # -------------------------------------------------
                result = st.session_state.conversational_chain.invoke({"input": prompt}, config={"configurable": {"session_id": session_id}})
                # Extract and clean all page_content from context documents
                if 'context' in result and result['context']:
                    all_page_contents = []
                    for index, doc in enumerate(result['context'], start=1):  # Start enumerating from 1
                        if hasattr(doc, 'page_content'):
                            cleaned_content = clean_page_content(doc.page_content)  # Clean the content
                            all_page_contents.append(f"Slide {index}: {cleaned_content}")  # Format for slides
                    
                    full_response = '\n'.join(all_page_contents) if all_page_contents else "No relevant information found."
                else:
                    full_response = "Sorry, I couldn't find any relevant information."

                st.session_state.chat_history.append({"user": prompt, "bot": full_response, "context_docs": []})

                # Button for summarizing the response
                # if st.button("Summarize Response"):
                #     summary = summarize_text(full_response)
                    # st.session_state.chat_history.append({"user": "Summary Request", "bot": summary, "context_docs": []})
                if 'context' in result and result['context']:
                    all_page_contents = []
                    for doc in result['context']:
                        if hasattr(doc, 'page_content'):
                            cleaned_content = clean_page_content(doc.page_content)  # Clean the content
                            all_page_contents.append(cleaned_content)
                    
                    response = ' '.join(all_page_contents) if all_page_contents else "No relevant information found."
                else:
                    response = "Sorry, I couldn't find any relevant information."

                st.session_state.chat_history.append({"user": prompt, "bot": response, "context_docs": []})

                # ------------------------------------
                # if 'context' in result and result['context']:
                #     page_content = result['context'][0].page_content  # Get the first document's page_content
                #     cleaned_content = clean_page_content(page_content)  # Clean the content
                #     response = page_content
                # else:
                #     response = "Sorry, I couldn't find any relevant information." 
                # session_id = "user412"
            #     response = st.session_state.conversational_chain.invoke({
            #     "input": prompt,
            #     "context": source_document  # Add context from the source document
            # }, config={"configurable": {"session_id": session_id}})
                # context_docs = response.get('context', [])
                # session_id = "user412"
                # st.session_state.chat_history.append({"user": prompt, "bot": response, "context_docs": []})
                # st.session_state.chat_history.insert(0, {"user": user_input, "bot": response['answer'], "context_docs": response.get("context_docs", [])})
                # response = st.session_state.conversational_chain.invoke({"input": prompt}, config={"configurable": {"session_id": session_id}})

                # response = st.session_state.conversational_chain.invoke({"input": prompt},
                # {"context":source_document}, 
                # # "context_docs": response.get("context_docs", [])} # Add context from the source document
                # config={"configurable": {"session_id": session_id}})

            #     bot_response = f"This is a response to the recommended prompt: '{response}'"
            #     st.markdown(f"**Bot:** {bot_response}")

            #     st.session_state["chat_history"].append({
            #     "user": prompt,
            #     "bot": bot_response,
            #     "timestamp": datetime.now().strftime("%Y-%m-%d %H:%M:%S")
            # })
#         def clean_text_response(response):
#             cleaned_texts = [doc.page_content for doc in response.get('context', [])]
#             return "\n".join(cleaned_texts)  # Combine all document contents into a single string
#         Document=""
# # Your original response
#         response = {
#             'input': 'Explain the logistics solutions Saksoft offer.',
#             'context': [
#             Document(metadata={'source': 'Data\\SalesDocs\\Saksoft- Top Case Studies by BUs.pdf', 'page': 11}, page_content='www.saksoft.com\n 12\nLogistics'),
#             Document(metadata={'source': 'Data\\SalesDocs\\Saksoft - Corporate Overview (1).pdf', 'page': 0}, page_content='Saksoft serves as a forward-thinking partner in digital transformation...')
#     ]
# }

# # Get the cleaned text response (without metadata)
#         cleaned_response = clean_text_response(response)

# # Output the clean text response
#         st.markdown(f"**Response:** {cleaned_response}")
            # iif msg.get("context_docs"):
            # with st.expander("Source Documents"):
            #     for doc in msg["context_docs"]:
            #         st.write(f"Source: {doc.metadata.get('source', 'N/A')}")
            #         st.write(doc.page_content)
    # if st.button("💡 Recommended Prompts"):
    #     st.session_state["active_section"] = "Recommended Prompts"
    # #     elif st.session_state["active_section"] == "Recommended Prompts":
    # # st.subheader("💡 Recommended Prompts")
    #     for prompt in st.session_state["recommended_prompts"]:
    #         if st.button(prompt):
    #         # Display selected prompt
    #             st.markdown(f"**You selected prompt:** {prompt}")
    #             response = st.session_state.conversational_chain.invoke({"input": prompt}, config={"configurable": {"session_id": session_id}})

    #             bot_response = f"This is a response to the recommended prompt: '{response}'"
    #             st.markdown(f"**Bot:** {bot_response}")
    #             st.session_state["chat_history"].append({
    #             "user": prompt,
    #             "bot": bot_response,
    #             "timestamp": datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    #         })
        

# Main content based on active section
st.markdown("""
<div style="background-color: #f0f2f5; padding: 10px; border-radius: 5px;">
<h3 style="color: #333;">Chat with Your Assistant</h3>
 </div>
""", unsafe_allow_html=True)           
# Sidebar for navigation
# st.sidebar.title("Chat Guide")
# st.sidebar.title("Assistant Portal")
# sidebar_option = st.sidebar.radio("Select an option:", ( "Contact Support"))


# if sidebar_option == "Help Page":
#     st.sidebar.write("Here you can find help regarding using the Sales Companion Bot.")
#     # Add more help content as needed

# elif sidebar_option == "View More":
#     st.sidebar.write("Additional resources and information can be added here.")
    # Add more content for "View More" as needed

# elif sidebar_option == "Contact Support":
#     st.sidebar.write("For support, please contact us at:")
#     st.sidebar.write("Email: support@example.com")
#     st.sidebar.write("Phone: +123 456 7890")

def generate_presentation(content):
    ppt_file_path = "generatedPPT1.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title = slide.shapes.title
    content_shape = slide.shapes.placeholders[1]

    title.text = "Generated Presentation"
    content_shape.text = content

    presentation.save(ppt_file_path)
    return ppt_file_path

if 'feedback' not in st.session_state:
    st.session_state.feedback = []
# Feedback mechanism
def handle_feedback(user_input):
    st.session_state.feedback.append(user_input)
    st.success("Thank you for your feedback!")
# Function to display feedback options
def display_feedback():
    st.markdown('<p style="font-weight:bold; color:black;">Was this response helpful?</p>', unsafe_allow_html=True)

    # st.write("Was this response helpful?")
    col1, col2 = st.columns(2)
    with col1:
        if st.button("👍 Like"):
            handle_feedback("Like")
    with col2:
        if st.button("👎 Dislike"):
            handle_feedback("Dislike")

    # Manual feedback option
    manual_feedback = st.text_area("Provide additional feedback (optional):")
    if st.button("Submit Feedback"):
        if manual_feedback:
            handle_feedback(manual_feedback)
# Custom CSS for button styling
st.markdown("""
    <style>
        .hunting-button {
            background-color: #ff4d4d;  /* Bold red color */
            color: white;
            font-size: 20px;
            font-weight: bold;
            padding: 10px 20px;
            border: none;
            border-radius: 5px;
            cursor: pointer;
            transition: all 0.3s;
            display: inline-flex;
            align-items: center;
        }
        .hunting-button:hover {
            background-color: #ff1a1a;  /* Darker red on hover */
            box-shadow: 0 0 10px rgba(255, 255, 255, 0.5);
            animation: pulsate 1s infinite;
        }
        @keyframes pulsate {
            0% { transform: scale(1); }
            50% { transform: scale(1.05); }
            100% { transform: scale(1); }
        }
        .hunting-icon {
            margin-right: 8px;
        }
    </style>
""", unsafe_allow_html=True)   
# Home page buttons for Hunting and Farming
st.write("Select your approach:")
# Recommended prompts
# recommended_prompts = [
#     "Summarize case studies for Data Analytics",
#     "Tell me the top 3 case studies",
#     "Explain the logistics solutions Saksoft offers"
# ]


# # Sidebar for recommended prompts with scrolling
# with st.sidebar:
#     st.markdown("<h4>Recommended Prompts</h4>", unsafe_allow_html=True)
    
#     # Create a scrolling area for the prompts
#     for prompt in recommended_prompts:
#         if st.button(prompt):
#             st.session_state.user_input = prompt
#             session_id = "user412"
#             response = st.session_state.conversational_chain.invoke({"input": prompt}, config={"configurable": {"session_id": session_id}})
            
#             # Get context documents
#             context_docs = response.get('context', [])
            
#             # Add the response to chat history
#             st.session_state.chat_history.append({
#                 "user": prompt,
#                 "bot": response.get('answer', "I didn't understand that."),
#                 "context_docs": context_docs
#             })

#     # Add some space to allow scrolling
#     st.markdown("<div style='height: 100px;'></div>", unsafe_allow_html=True)  # Adjust height as needed
# Recommended prompts
# recommended_prompts = [
    # "Summarize case studies for Data Analytics",
    # "Tell me the top 3 case studies",
    # "Explain the logistics solutions Saksoft offers"  # Ensure this prompt is correctly spelled and formatted
# ]

# # Sidebar for recommended prompts
# with st.sidebar:
#     st.markdown("<h4>Recommended Prompts</h4>", unsafe_allow_html=True)
    
#     for prompt in recommended_prompts:
#         if st.button(prompt):
#             st.session_state.user_input = prompt
            
#             # Ensure session_id is defined
#             session_id = "user412"
            
#             # Prepare the input for the conversational chain
#             input_data = {
#                 "chat_history": st.session_state['chat_history'],  # Include chat history
#                 "question": prompt  # Include the current prompt as the question
#             }
            
#             try:
#                 # Invoke the conversational chain
#                 response = st.session_state.conversational_chain.invoke(input_data, config={"configurable": {"session_id": session_id}})
                
#                 # Check if response contains expected keys
#                 if 'answer' in response:
#                     context_docs = response.get('context', [])
#                     st.session_state['chat_history'].append({
#                         "user": prompt,
#                         "bot": response['answer'],
#                         "context_docs": context_docs
#                     })
#                 else:
#                     # Handle unexpected response structure
#                     st.session_state['chat_history'].append({
#                         "user": prompt,
#                         "bot": "I'm sorry, but I didn't get a valid response.",
#                         "context_docs": []
#                     })
                    
#             except Exception as e:
#                 # Handle any exceptions that occur during invocation
#                 st.session_state['chat_history'].append({
#                     "user": prompt,
#                     "bot": f"An error occurred: {str(e)}",
#                     "context_docs": []
#                 })
# ------------------------below is main code------------------
# recommended_prompts = [
#     "Summarize case studies for Data Analytics",
#     "Tell me the top 3 case studies",
# ]

# Sidebar for recommended prompts
# with st.sidebar:
#     st.markdown("<div class='box'><h4>Recommended Prompts</h4></div>", unsafe_allow_html=True)
#     for prompt in recommended_prompts:
#         if st.button(prompt):
#             st.session_state.user_input = prompt
#             # Simulating a response for the clicked prompt
#             if prompt == "Summarize case studies for Data Analytics":
#                 # response = {"answer": "Here are some key insights from the case studies in Data Analytics...", "context_docs": []}
#                 session_id ="user412"
#                 response = st.session_state.conversational_chain.invoke({"input": prompt}, config={"configurable": {"session_id": session_id}})
#                 context_docs = response.get('context', [])
#             elif prompt == "Tell me the top 3 case studies":
#                 # response = {"answer": "1. Case Study A\n2. Case Study B\n3. Case Study C", "context_docs": []}
#                 session_id ="user412"
#                 response = st.session_state.conversational_chain.invoke({"input": prompt}, config={"configurable": {"session_id": session_id}})
#                 context_docs = response.get('context', [])
#             elif prompt == "Explain the logistics solution saksof offer":
#                 # response = {"answer": "1. Case Study A\n2. Case Study B\n3. Case Study C", "context_docs": []}
#                 session_id ="user412"
#                 response = st.session_state.conversational_chain.invoke({"input": prompt}, config={"configurable": {"session_id": session_id}})
#                 context_docs = response.get('context', [])
            
            # Add the response to chat history
            # st.session_state.chat_history.append({
            #     "user": prompt,
            #     # "bot": response.get('answer', "I didn't understand that."),
            #     "context_docs": response.get("context_docs", [])
            # })
# Buttons for Hunting and Farming
col1, col2,col3 = st.columns(3)

with col1:
    if st.button("Hunting"):
        st.session_state.mode = "hunting"
        st.write("You selected Hunting mode.")

with col2:
    if st.button("Farming"):
        st.session_state.mode = "farming"
        st.write("You selected Farming mode.")
with col3:
    if st.button("Refresh"):
        st.session_state.user_input = ""  # Clear the input
        st.session_state.chat_history = []         
        # Function to generate PowerPoint (Placeholder)
# Function to generate PowerPoint (Placeholder)
def generate_presentation(content):
    # Here you would implement the logic to create a PPT based on the content
    # This is a placeholder function that simulates PPT generation
    ppt_file_path = "generated_presentation1.pptx"  # Specify your file path
    # Create a PPT file and save it (this is a simplified example)
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title = slide.shapes.title
    content_shape = slide.shapes.placeholders[1]

    title.text = "Generated Presentation"
    content_shape.text = content

    presentation.save(ppt_file_path)
    return ppt_file_path

def generate_ppt(response_text, file_name="response_presentation.pptx"):
    """
    Generate a PowerPoint file dynamically based on the response text.
    """
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "Query Response"
    content.text = response_text

    prs.save(file_name)
    return file_name
# Feedback handling
def capture_feedback(query, feedback_type, additional_feedback=None):
    st.session_state.feedback.append({
        "query": query,
        "feedback_type": feedback_type,
        "additional_feedback": additional_feedback
    })
    st.success("Thank you for your feedback!")

# Display a textbox for prompts when Farming mode is selected
if st.session_state.get('mode') == "farming":
    st.header("Farming Approach")
    st.markdown("""
    The farming approach in sales often involves personalized attention to clients, understanding their needs over time, and ensuring that the solutions offered evolve with their changing demands.
    """)
    user_input = st.text_input("Enter your prompt for Farming:")

    if st.button("Submit"):
        user_input_lower = user_input.lower()  # Normalize input for checking

        # Store the user input in session state
        st.session_state.user_input = user_input

        # Respond to greetings
        if any(greet in user_input_lower for greet in ["hi", "hello", "hey"]):
            response = "Hello! How can I assist you?"
            st.session_state.chat_history.append({"user": user_input, "bot": response, "context_docs": []})
        # Respond to thanks
        elif any(thank in user_input_lower for thank in ["thank you", "thanks", "thank you very much", "thx"]):
            response = "Thank you! Let me know if you have any more queries."
            st.session_state.chat_history.append({"user": user_input, "bot": response, "context_docs": []})
        # Respond to yes/no
        elif any(yn in user_input_lower for yn in ["yes", "no"]):
            response = "I understand. If you have any other questions, feel free to ask!"
        
        # Process document queries
        elif user_input and 'conversational_chain' in st.session_state:
            session_id = "user412"  # Static session ID for this demo; you can make it dynamic if needed
            response = st.session_state.conversational_chain.invoke({"input": user_input}, config={"configurable": {"session_id": session_id}})
            if "meeting a new prospect" in user_input_lower or "latest corporate presentation" in user_input_lower:
                response = {
                    "answer": "Please find below the document for your query.",
                    "context_docs": split_docs  # Show all documents related to the query
                }

        # else:
        #         session_id ="user412"
        #         response = st.session_state.conversational_chain.invoke({"input": user_input_lower}, config={"configurable": {"session_id": session_id}})
        #         context_docs = response.get('context', [])
        # st.session_state.chat_history.insert(0, {"user": user_input, "bot": response['answer'], "context_docs": response.get("context_docs", [])})

        # # #  Display feedback options
        # display_feedback()
        #          # Insert the new response at the top of the chat history
        # session_id ="user412"
        # response = st.session_state.conversational_chain.invoke({"input": user_input}, config={"configurable": {"session_id": session_id}})

        # st.session_state.chat_history.insert(0, {"user": user_input, "bot": response['answer'], "context_docs": response.get("context_docs", [])})
                # st.divider()  # Visual divider
                # st.markdown("##### Was this response helpful?")
                # col1, col2 = st.columns(2)

# Feedback buttons
                # with col1:
                #  if st.button("👍 Like"):
                #    capture_feedback(user_input, "like")

                # with col2:
                #   if st.button("👎 Dislike"):
                #     additional_feedback = st.text_input("If you'd like, please provide additional feedback:")
                #     if st.button("Submit Feedback"):
                #         capture_feedback(user_input, "dislike", additional_feedback)
        session_id ="user412"       
        response = st.session_state.conversational_chain.invoke({"input": user_input}, config={"configurable": {"session_id": session_id}})
        context_docs = response.get('context', [])
        st.session_state.chat_history.insert(0, {"user": user_input, "bot": response['answer'], "context_docs": response.get("context_docs", [])})

        #  Display feedback options
        display_feedback()
        #response = st.session_state.conversational_chain.invoke({"input": user_input_hunting}, config={"configur
# Display a textbox for prompts when Hunting mode is selecte
path="Data\corporateoverview"
import fitz  # PyMuPDF

# def extract_summary_from_pdf(path):
#     summary = ""
#     try:
#         with fitz.open(path) as pdf_document:
#             for page_num in range(len(pdf_document)):
#                 page = pdf_document[page_num]
#                 summary += page.get_text()
        
#         return summary[:]  # Return the first 1000 characters for brevity
#     except Exception as e:
#         return f"An error occurred: {e}"
# Helper function to extract text from PDF
# def extract_pdf_text(path):
#     # try:
#         reader = PdfReader(path)
        
#         # Check if the PDF has pages
#         if len(reader.pages) == 0:
#             return "Error: The PDF file is empty."

#         raw_text = ""
#         for page in reader.pages:
#             text = page.extract_text()
#             if text:  # Only append if text extraction was successful
#                 raw_text += text + " "
        
#         # Clean the text: remove unwanted symbols, excessive spaces, and line breaks
#         clean_text = re.sub(r'\s+', ' ', raw_text)  # Replace multiple spaces/newlines with a single space
#         clean_text = re.sub(r'[^\w\s.,:;!?-]', '', clean_text)  # Remove special characters (except punctuation)

#         # Split into sentences for better readability
#         sentences = re.split(r'(?<=[.?!])\s+', clean_text)

#         # Arrange the sentences in a point-to-point manner
#         structured_text = "\n".join([f"- {sentence.strip()}" for sentence in sentences if sentence.strip()])

#         # Summarize the structured text
#         summarizer = pipeline("summarization")
#         summary = summarizer(structured_text, max_length=100, min_length=30, do_sample=False)

        # Extract the summarized text
        # summarized_text = summary[0]['summary_text']

        # return summarized_text.strip()
def extract_pdf_text(pdf_path):
    # try:
        reader = PdfReader(pdf_path)
        raw_text = " ".join([page.extract_text() for page in reader.pages])
        # Clean the text: remove unwanted symbols, excessive spaces, and line breaks
        clean_text = re.sub(r'\s+', ' ', raw_text)  # Replace multiple spaces/newlines with a single space
        clean_text = re.sub(r'[^\w\s.,:;!?-]', '', clean_text)  # Remove special characters (except punctuation)
 
        # Split into sentences for better readability
        sentences = re.split(r'(?<=[.?!])\s+', clean_text)
 
        # Arrange the sentences in a point-to-point manner
        structured_text = "\n".join([f"- {sentence.strip()}" for sentence in sentences if sentence.strip()])
 
        return structured_text
#     except Exception as e:
#         return f"Error reading or processing the PDF: {str(e)}"
    # except Exception as e:
    #     return f"Error reading or processing the PDF: {str(e)}"

# Example usage:
    # pdf_path = "Data/corporateoverview/Saksoft Overview_Nov24.pdf"
    # cleaned_text = extract_pdf_text(pdf_path)
    # print(cleaned_text)

    #     return document_text
    # except Exception as e:
    #     return f"Error reading PDF: {str(e)}"
if st.session_state.get('mode') == "hunting":
    st.header("Hunting Approach")
    st.markdown("""
    The hunting approach in sales focuses on acquiring new customers and expanding market reach. It involves identifying potential clients, understanding their needs, and offering solutions that address their pain points.
    """)
    user_input_hunting = st.text_input("Enter your prompt for Hunting:")
    path="Data\corporateoverview\Saksoft Overview_Nov24.pdf"
    url="https://saksoftonline.sharepoint.com/:p:/s/CopilotBOT/EYvJGu8IY0RHvPnOHvqI3pEBU7C6nRDLHEbznpUTJtw7nw?e=dFTmvo"
    if st.button("Submit Hunting"):
        user_input_hunting_lower = user_input_hunting.lower()
        st.session_state.user_input_hunting = user_input_hunting
        path="Data\corporateoverview\Saksoft Overview_Nov24.pdf"
              # Check for specific queries
        # if "meeting a new prospect" in user_input_hunting or "latest corporate presentation" or "Give me summary of saksoft corporate overview" in user_input_hunting:

# Handle user input
        if "i am meeting a new prospect" in user_input_hunting.lower() :
            # try:
        # Extract text from the PDF
            pdf_text = extract_pdf_text(path)

            if "Error" in pdf_text:
                st.error(pdf_text)  # Show error if PDF extraction fails
            else:
            # Generate a summary (example: first 1500 characters)
                summary = pdf_text[:500] + "..." if len(pdf_text) > 1500 else pdf_text
                # summary = pdf_text

            # Display the summary
                st.write("### Saksoft Corporate Overview Summary")
                st.write(summary)
                pdf_link = os.path.abspath(path)
            # st.markdown(f"[📄 Open and Download the PDF]({url})", unsafe_allow_html=True)
            st.markdown(
                f'<a href="{url}" target="_blank" style="text-decoration:none;color:blue;font-weight:bold;">📄 Open and Download the Full PDF</a>',
            unsafe_allow_html=True,
            )
        elif "give me summary of corporate overview" in user_input_hunting.lower():
            # Extract text from the PDF
            url="https://saksoftonline.sharepoint.com/:p:/s/CopilotBOT/EYvJGu8IY0RHvPnOHvqI3pEBU7C6nRDLHEbznpUTJtw7nw?e=dFTmvo"
            pdf_text = extract_pdf_text(path)


            if "Error" in pdf_text:
                st.error(pdf_text)  # Show error if PDF extraction fails
            else:
            # Generate a summary (example: first 1500 characters)
                summary = pdf_text[:1500] + "..." if len(pdf_text) > 1500 else pdf_text

            # Display the summary
            st.write("### Saksoft Corporate Overview Summary")
            st.write(summary)
            pdf_link = os.path.abspath(path)
            # st.markdown(f"[📄 Open and Download the PDF]({url})", unsafe_allow_html=True)
            st.markdown(
                f'<a href="{url}" target="_blank" style="text-decoration:none;color:blue;font-weight:bold;">📄 Open and Download the Full PDF</a>',
            unsafe_allow_html=True,
            )
        elif any(qu in user_input_hunting for qu in ["quantum mechanics", "facebook"]):
            response = "Sorry, could not find relevant information in documents. Please rephrase your query."
        elif  "quantum mechanics" in user_input_hunting.lower():
        # elif any(yn in user_input_lower for yn in ["yes", "no"]):
            response = "Sorry, could not find relevant information in documents. Please rephrase your query."    
        elif  "facebook" in user_input_hunting.lower():
        # elif any(yn in user_input_lower for yn in ["yes", "no"]):
            response = "Sorry, could not find relevant information in documents. Please rephrase your query." 
   
        else:
          session_id ="user412"
          response = st.session_state.conversational_chain.invoke({"input": user_input_hunting}, config={"configurable": {"session_id": session_id}})
          response_text = response['answer']


           # Insert the new response at the top of the chat history
        #   st.session_state.chat_history.insert(0, {"user": user_input_hunting, "bot": response['answer'], "context_docs": response.get("context_docs", [])})

#     session_id ="user412"
# response = st.session_state.conversational_chain.invoke({"input": user_input_hunting}, config={"configurable": {"session_id": session_id}})
# Generate a PPT dynamically for the response
        #   ppt_file = generate_ppt(response_text)
        #   st.markdown(f"**Bot Response:** {response_text}")
        #   with open(ppt_file, "rb") as f:
        #       st.download_button(
        #       label="Download Generated PPT",
        #       data=f,
        #       file_name=ppt_file,
        #       mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        # )
        session_id ="user412"      
        response = st.session_state.conversational_chain.invoke({"input": user_input_hunting}, config={"configurable": {"session_id": session_id}})
        # response =summary
        st.session_state.chat_history.insert(0, {"user": user_input_hunting, "bot": response['answer'], "context_docs": response.get("context_docs", [])})
#         st.divider()  # Visual divider
#         st.markdown("##### Was this response helpful?")
#         col1, col2 = st.columns(2)

# # Feedback buttons
#         with col1:
#             if st.button("👍 Like"):
#                 capture_feedback(user_input_hunting, "like")

#         with col2:
#             if st.button("👎 Dislike"):
#                 additional_feedback = st.text_input("If you'd like, please provide additional feedback:")
#                 if st.button("Submit Feedback"):
#                     capture_feedback(user_input_hunting, "dislike", additional_feedback)
  
        # st.session_state.chat_history.insert(0, {"user": user_input_hunting, "bot": response['answer'], "context_docs": response.get("context_docs", [])})
        # Display feedback options
        display_feedback()
  # else:


#     session_id = "user412"
#     response = st.session_state.conversational_chain.invoke({"input": user_input}, config={"configurable": {"session_id": session_id}})
#     context_docs = response.get('context', [])
      
        # if "i am meeting a new prospect" in user_input_hunting_lower or "latest corporate presentation" in user_input_hunting_lower:
        #     ppt_file = generate_presentation("This presentation covers the latest corporate strategies and insights.")
        #     st.write("Please find below the document for your query.")
        #     st.download_button("Download Presentation", ppt_file, file_name="presentation.pptx")
        # Check for specific prompts and generate PPT
        # if "meeting a new prospect" in user_input_hunting_lower:
        #     ppt_response = generate_presentation("Help me with the latest corporate presentation.")
        #     st.write(ppt_response)


#  elif "top customers for" in user_input_hunting_lower:
#             ppt_response = generate_presentation("Help with top customers for [Industry], [Technology].")
#             st.write(ppt_response)    
#         else:
#                 session_id ="user412"
#                 response = st.session_state.conversational_chain.invoke({"input": user_input_hunting}, config={"configurable": {"session_id": session_id}})
#                 context_docs = response.get('context', [])
#         #          # Insert the new response at the top of the chat history
#         session_id ="user412"
#         response = st.session_state.conversational_chain.invoke({"input": user_input_hunting}, config={"configurable": {"session_id": session_id}})

#         st.session_state.chat_history.insert(0, {"user": user_input_hunting, "bot": response['answer'], "context_docs": response.get("context_docs", [])})
           


        #response = st.session_state.conversational_chain.invoke({"input": user_input_hunting}, config={"configurable": {"session_id": session_id}})
        # if isinstance(response, str):
        #             # If response is a string, wrap it in a dictionary
        #             response = {"answer": response, "context_docs": []}
# Display chat history with responses
if st.session_state.get("chat_history", []):
    for msg in st.session_state["chat_history"]:
        st.markdown(f"**You:** {msg['user']}", unsafe_allow_html=True)

        # Check for specific prompts
        if "i am meeting a new prospect" in msg["user"].lower():
            pdf_text = extract_pdf_text(path)
            if "Error" in pdf_text:
                bot_response = f"An error occurred: {pdf_text}"
            else:
                # Generate a summary (example: first 1500 characters)
                url="https://saksoftonline.sharepoint.com/:p:/s/CopilotBOT/EYvJGu8IY0RHvPnOHvqI3pEBU7C6nRDLHEbznpUTJtw7nw?e=dFTmvo"

                summary = pdf_text[:1500] + "..." if len(pdf_text) > 1500 else pdf_text
                bot_response = summary
                pdf_link = os.path.abspath(path)
                # st.markdown(f"[📄 Open and Download the Full PDF](file://{pdf_link})", unsafe_allow_html=True)
                st.markdown(
                    f'<a href="{url}" target="_blank" style="text-decoration:none;color:blue;font-weight:bold;">📄 Open and Download the Full PDF</a>',
                    unsafe_allow_html=True,)
        elif  "give me summary of corporate overview" in msg["user"].lower():
            pdf_text = extract_pdf_text(path)
            if "Error" in pdf_text:
                bot_response = f"An error occurred: {pdf_text}"
            else:
                # Generate a summary (example: first 1500 characters)
                url="https://saksoftonline.sharepoint.com/:p:/s/CopilotBOT/EYvJGu8IY0RHvPnOHvqI3pEBU7C6nRDLHEbznpUTJtw7nw?e=dFTmvo"

                summary = pdf_text[:1500] + "..." if len(pdf_text) > 1500 else pdf_text
                bot_response = summary
                pdf_link = os.path.abspath(path)
                # st.markdown(f"[📄 Open and Download the Full PDF](file://{pdf_link})", unsafe_allow_html=True)
                st.markdown(
                    f'<a href="{url}" target="_blank" style="text-decoration:none;color:blue;font-weight:bold;">📄 Open and Download the Full PDF</a>',
                    unsafe_allow_html=True,)
        
        else:
            st.markdown(f"**Bot:** {msg['bot']}", unsafe_allow_html=True)

# with col3:
#     if st.button("Refresh"):
#         st.session_state.user_input = ""  # Clear the input
#         st.session_state.chat_history = [] 
        # else:
            # General conversational_chain response
            # bot_response = "The conversational_chain generates a response based on the provided documents."

        # st.markdown(f"**Bot:** {msg['bot']}", unsafe_allow_html=True)
# ==================working below=====================
# if st.session_state["chat_history"]:
#     for msg in st.session_state["chat_history"]:
#         st.markdown(f"**You:** {msg['user']}",unsafe_allow_html=True)
#         st.markdown(f"**Bot:** {msg['bot']}", unsafe_allow_html=True)
                
#         if msg.get("context_docs"):
#             with st.expander("Source Documents"):
#                 for doc in msg["context_docs"]:
#                     st.write(f"Source: {doc.metadata.get('source', 'N/A')}")
#                     st.write(doc.page_content)
# -------------------------------
    #  Display feedback options
    #     display_feedback()

# with col3:
#     if st.button("Refresh"):
#         st.session_state.user_input = ""  # Clear the input
#         st.session_state.chat_history = [] 


        # # Storing the response in chat history
        # st.session_state.chat_history.append({
        #     "user": user_input,
        #     # "bot": response.get('answer', "I didn't understand that."),
        #     "context_docs": response.get("context_docs", [])
        # })
# if st.session_state["chat_history"]:
#     for msg in st.session_state["chat_history"]:
#         st.markdown(f"**You:** {msg['user']}")
#         st.markdown(f"**Bot:** {msg['bot']}")



# if st.session_state["chat_history"]:
#     for msg in st.session_state["chat_history"]:
#         st.markdown(f"**You:** {msg['user']}",unsafe_allow_html=True)
#         st.markdown(f"**Bot:** {msg['bot']}", unsafe_allow_html=True)
#         if msg.get("context_docs"):
#             with st.expander("Source Documents"):
#                 for doc in msg["context_docs"]:
#                     st.write(f"Source: {doc.metadata.get('source', 'N/A')}")
#                     st.write(doc.page_content) 




# if st.session_state.chat_history:
#     for message in st.session_state.chat_history:
#         st.markdown(user_template.format(msg=message['user']), unsafe_allow_html=True)
#         st.markdown(bot_template.format(msg=message['bot']), unsafe_allow_html=True)

#         if message.get("context_docs"):
#             with st.expander("Source Documents"):
#                 for doc in message["context_docs"]:
#                     st.write(f"Source: {doc.metadata.get('source', 'N/A')}")
#                     st.write(doc.page_content)       
# # Chat history display
# if st.session_state.chat_history:
#     chat_container = st.container()
#     with chat_container:
#         for message in st.session_state.chat_history:
#             st.markdown(f"<div class='user-message'>{message['user']}</div>", unsafe_allow_html=True)
#             st.markdown(f"<div class='bot-message'>{message['bot']}</div>", unsafe_allow_html=True)

# # Chat history display
# if st.session_state.chat_history:
#     chat_container = st.container()
#     with chat_container:
#         for message in st.session_state.chat_history:
#             st.markdown(f"<div class='user-message'>{message['user']}</div>", unsafe_allow_html=True)
            # st.markdown(f"<div class='bot-message'>{message['bot']}</div>", unsafe_allow_html=True)

# # Sidebar for frequently asked questions
# with st.sidebar:
#     st.subheader("Frequently Asked Questions")
#     st.markdown("<div class='sidebar-content'>", unsafe_allow_html=True)  # Sidebar content styling
#     for message in st.session_state.chat_history:
#         if st.button(message['user']):
#             st.session_state.selected_question = message['user']
#             st.session_state.selected_answer = message['bot']
#     st.markdown("</div>", unsafe_allow_html=True)  # Close sidebar content styling
# st.image(logopath, width=200)
# st.title("Sales Companion Bot")
# st.write("Welcome! How can I help you today?")

# pdf_directory = "Data/SalesDocs"  # Update to your actual folder path

# # Prepare documents and ingest into vector store
# split_docs = prepare_and_split_docs(pdf_directory)
# vector_db = ingest_into_vectordb(split_docs)
# retriever = vector_db.as_retriever()

# # Initialize conversation chain
# conversational_chain = get_conversation_chain(retriever)
# st.session_state.conversational_chain = conversational_chain

# if 'chat_history' not in st.session_state:
#     st.session_state.chat_history = []

# # Chat input
# # conversation_chain = initialize_conversation_chain(llm)
# user_input = st.text_input("Ask a question about the provided documents:")
#

# Chat history display
# if st.session_state.chat_history:
#     chat_container = st.container()
#     with chat_container:
#         for message in st.session_state.chat_history:
#             st.markdown(f"<div class='user-message'>{message['user']}</div>", unsafe_allow_html=True)
#             st.markdown(f"<div class='bot-message'>{message['bot']}</div>", unsafe_allow_html=True)

# Sidebar for frequently asked questions
# with st.sidebar:
    # st.subheader("Frequently Asked Questions")
    # st.markdown("<div class='sidebar-content'>", unsafe_allow_html=True)  # Sidebar content styling
    # for message in st.session_state.chat_history:
    #     if st.button(message['user']):
    #         st.session_state.selected_question = message['user']
    #         st.session_state.selected_answer = message['bot']
    # st.markdown("</div>", unsafe_allow_html=True)  # Close sidebar content styling

# col1, col2 = st.columns([1, 1])
# +++++++++++++++++++
# with col1:
#     if st.button("Submit"):
#         user_input_lower = user_input.lower()  # Normalize input for checking
        

#     #     # Store the user input in session state
#     #     st.session_state.user_input = user_input
#     #     # Respond to greetings
#         if any(greet in user_input_lower for greet in ["hi", "hello", "hey"]):
#             response = "Hello! How can I assist you?"
#             st.session_state.chat_history.append({"user": user_input, "bot": response, "context_docs": []})

#         # Respond to thanks
#         elif any(thank in user_input_lower for thank in ["thank you", "thanks", "thank you very much", "thx"]):
#             response = "Thank you! Let me know if you have any more queries."
#             st.session_state.chat_history.append({"user": user_input, "bot": response, "context_docs": []})

#         # Respond to yes/no
#         elif any(yn in user_input_lower for yn in ["yes", "no"]):
#             response = "I understand. If you have any other questions, feel free to ask!"
#             st.session_state.chat_history.append({"user": user_input, "bot": response, "context_docs": []})

#         # Process document queries
#         elif user_input and 'conversational_chain' in st.session_state:
#             session_id ="user412" # Static session ID for this demo; you can make it dynamic if needed
#             response = st.session_state.conversational_chain.invoke({"input": user_input}, config={"configurable": {"session_id": session_id}})
#             # context_docs = response.get('context', [])


#         # Process document queries
#         # elif user_input and 'conversational_chain' in st.session_state:
#         #     session_id = "user412"  # Static session ID for this demo; you can make it dynamic if needed
#         #     response = st.session_state.conversational_chain.invoke({"input": user_input}, config={"configurable": {"session_id": session_id}})

        #     # Check for specific queries
        # if "meeting a new prospect" in user_input_lower or "latest corporate presentation" in user_input_lower:
        #         response = {
        #             "answer": "Please find below the document for your query.",
        #             "context_docs": split_docs  # Assuming you want to show all documents related to the query
        #         }
#         else:
#                 session_id ="user412"
#                 response = st.session_state.conversational_chain.invoke({"input": user_input}, config={"configurable": {"session_id": session_id}})
#                 context_docs = response.get('context', [])
#         #          # Insert the new response at the top of the chat history
#         st.session_state.chat_history.insert(0, {"user": user_input, "bot": response['answer'], "context_docs": response.get("context_docs", [])})
#         # Check for specific queries
#         # Check for specific queries
#         session_id ="user412"
#         if "meeting a new prospect" in user_input_lower or "latest corporate presentation" in user_input_lower:
#             session_id ="user412"
#     # Create a response with specific document names and paths
#             response = {
#         "answer": "Please find below the document for your query.",
#         "context_docs": [
#             {"name": "Meeting_New_Prospect.pdf", "path": f"{pdf_directory}/Meeting_New_Prospect.pdf"},
#             {"name": "Latest_Corporate_Presentation.pdf", "path": f"{pdf_directory}/Latest_Corporate_Presentation.pdf"}
#         ]  # Only include relevant documents
#     }
#     else:
#         session_id ="user412"
#         response = st.session_state.conversational_chain.invoke({"input": user_input}, config={"configurable": {"session_id": session_id}})
# ++++++++++++++++++++++++/


            # Append the chat history
        # st.session_state.chat_history.append({"user": user_input, "bot": response['answer'], "context_docs": response.get("context_docs", [])})

# col1, col2 = st.columns([1, 1])
# with col1:
#     if st.button("Submit"):
#         user_input_lower = user_input.lower()  # Normalize input for checking

#         # Store the user input in session state
#         st.session_state.user_input = user_input
        

#         # Respond to greetings
#         if any(greet in user_input_lower for greet in ["hi", "hello", "hey"]):
#             response = "Hello! How can I assist you?"
#             st.session_state.chat_history.append({"user": user_input, "bot": response, "context_docs": []})

#         # Respond to thanks
#         elif any(thank in user_input_lower for thank in ["thank you", "thanks", "thank you very much", "thx"]):
#             response = "Thank you! Let me know if you have any more queries."
#             st.session_state.chat_history.append({"user": user_input, "bot": response, "context_docs": []})

#         # Respond to yes/no
#         elif any(yn in user_input_lower for yn in ["yes", "no"]):
#             response = "I understand. If you have any other questions, feel free to ask!"
#             st.session_state.chat_history.append({"user": user_input, "bot": response, "context_docs": []})

#         # Process document queries
#         elif user_input and 'conversational_chain' in st.session_state:
#             session_id ="user412" # Static session ID for this demo; you can make it dynamic if needed
#             response = st.session_state.conversational_chain.invoke({"input": user_input}, config={"configurable": {"session_id": session_id}})
#             context_docs = response.get('context', [])

#             # Check if the response is the default message
#             if response['answer'] == "I'm sorry, but I couldn’t find an answer to your question. Could you rephrase or provide more details?":
#                 st.session_state.chat_history.append({"user": user_input, "bot": response['answer'], "context_docs": []})
#                 # Check for specific queries
#                 if "meeting a new prospect" in user_input_lower or "latest corporate presentation" in user_input_lower:
#                  response = {
#                     "answer": "Please find below the document for your query.",
#                     "context_docs": split_docs  # Assuming you want to show all documents related to the query
#                 }
#             else:
#                 response = st.session_state.conversational_chain.invoke({"input": user_input}, config={"configurable": {"session_id": session_id}})
#                 context_docs = response.get('context', [])

#                 st.session_state.chat_history.append({"user": user_input, "bot": response['answer'], "context_docs": context_docs})
#                 st.session_state.chat_history.insert(0, {"role": "assistant", "content": context_docs})
#         # Use LangSmith for logging
#         # with langsmith.log_session("SalesBot"):
#         #     response = conversational_chain.invoke({"input": user_input}, config={"session_id": "user123"})
#         #     context_docs = response.get("context", [])
#         #     response = response["answer"]
#         #     if response == "I'm sorry, but I couldn’t find an answer.":
#         #         context_docs = []

#         # # Log response and context to LangSmith
#         # langsmith.log_event("Response", data={"response": response})
#         # langsmith.log_event("Context Docs", data={"docs": context_docs})

#     # st.session_state.chat_history.append({"user": user_input, "bot": response, "context_docs": context_docs})


# with col2:
#     if st.button("Refresh"):
#         st.session_state.user_input = ""  # Clear the input
#         st.session_state.chat_history = []  # Clear chat 
#     # if st.button("Download PPT"):
#         if st.session_state.chat_history:
        #     create_ppt(st.session_state.chat_history)
        #     st.success("PPT created! You can download it [here](chatbot_responses.pptx).")
        # else:
        #     st.warning("No chat history available to create a PPT.")
# with col2:
#     if st.button("Refresh"):
#         st.session_state.user_input = ""  # Clear the input

# col1, col2 = st.columns([1, 1])
# if st.button("Submit"):
#     user_input_lower = user_input.lower()
#     greetings = ["hi", "hello", "hey"]
#     thanks = ["thank you", "thanks", "thx"]

#     if any(greet in user_input_lower for greet in greetings):
#         response = "Hello! How can I assist you?"
#     elif any(thank in user_input_lower for thank in thanks):
#         response = "Thank you! Let me know if you have more queries."
#     else:
#         response = conversational_chain.invoke({"input": user_input}, config={"session_id": "user123"})
#         context_docs = response.get("context", [])
#         response = response["answer"]
#         if response == "I'm sorry, but I couldn’t find an answer.":
#             context_docs = []

#     st.session_state.chat_history.append({"user": user_input, "bot": response, "context_docs": context_docs})
# with col2:
#     if st.button("Next Question"):
#         st.session_state.user_input = ""
# response=""
# if st.session_state.chat_history:
#  for message in st.session_state.chat_history:
# # for message in reversed(st.session_state.chat_history):
#     user_msg = message.get("user")
#     bot_msg = message.get("bot")

    # st.markdown(user_template.format(msg=user_msg), unsafe_allow_html=True)
    # st.markdown(bot_template.format(msg=bot_msg), unsafe_allow_html=True)
# Display chat history
# if st.session_state.chat_history:
#     for message in st.session_state.chat_history:
#         st.markdown(user_template.format(msg=message['user']), unsafe_allow_html=True)
#         st.markdown(bot_template.format(msg=message['bot']), unsafe_allow_html=True)
# # for message in st.session_state.chat_history:
#     st.markdown(user_template.format(msg=message['user']), unsafe_allow_html=True)
#     st.markdown(bot_template.format(msg=message['bot']), unsafe_allow_html=True)

# ++++++++++++++++++++++++
# Display chat history/
# if st.session_state.chat_history:
#     for message in st.session_state.chat_history:
#         st.markdown(user_template.format(msg=message['user']), unsafe_allow_html=True)
#         st.markdown(bot_template.format(msg=message['bot']), unsafe_allow_html=True)

#         if message.get("context_docs"):
#             with st.expander("Source Documents"):
#                 for doc in message["context_docs"]:
#                     st.write(f"Source: {doc.metadata.get('source', 'N/A')}")
#                     st.write(doc.page_content)
# if st.session_state.chat_history:
#     for message in st.session_state.chat_history:
#         if 'user' in message and 'bot' in message:
#             st.markdown(user_template.format(msg=message['user']), unsafe_allow_html=True)
#             st.markdown(bot_template.format(msg=message['bot']), unsafe_allow_html=True)
#         else:
#             st.error("Invalid message format: missing 'user' or 'bot' key.")


#     if message.get("context_docs"):
#         with st.expander("Source Documents"):
#             for doc in message["context_docs"]:
#                 st.write(f"Source: {doc.metadata.get('source', 'N/A')}")
#                 st.write(doc.page_content)


# Display chat history (newest at the top)
    # st.write("### Chat History")
    # for message in st.session_state.messages:
    #     if message["role"] == "user":
    #         st.markdown(f"**You:** {message['content']}")
    #     else:
    #         st.markdown(f"**Bot:** {message['content']}")

# if __name__ == "__main__":
#     main()

# def display_chat_history(session_id):
#     history = store.get(session_id, ChatMessageHistory())
#     for msg in history.messages:
#         if msg.type == "human":
#             st.write(f"**You:** {msg.content}")
#         elif msg.type == "system":
#             st.write(f"**Bot:** {msg.content}")
# +++++++++++++++++++++++++++++++++++++++++++++++
# import streamlit as st
# # from your_module import prepare_and_split_docs, ingest_into_vectordb, get_conversation_chain

# # Display the logo and welcome message
# logopath = r"C:\Users\Akshada.P\OneDrive - SAKSOFT LIMITED\Documents\SAKSOFT Internal BOT\saksoftimg.png"
# st.image(logopath, width=200)
# st.title("Q&A Bot")
# st.write("Welcome! How can I help you today?")

# # Define the PDF directory (update path as per your environment)
# pdf_directory = "Data/SalesDocs"

# # Prepare documents and ingest them into vector storage
# try:
#     st.write("Preparing documents for search...")
#     split_docs = prepare_and_split_docs(pdf_directory)  # Function to read and split PDFs
#     vector_db = ingest_into_vectordb(split_docs)       # Function to ingest documents into a vector DB
#     retriever = vector_db.as_retriever()               # Retrieve documents using a retriever interface

#     # Initialize conversation chain
#     conversational_chain = get_conversation_chain(retriever)
#     st.session_state.conversational_chain = conversational_chain

#     # Initialize chat history if not already in session state
#     if 'chat_history' not in st.session_state:
#         st.session_state.chat_history = []
# except Exception as e:
#     st.error(f"Error preparing documents: {e}")

# # Chat Input Section
# user_input = st.text_input("Ask a question about the documents:")

# # Buttons for submitting and managing chat
# col1, col2 = st.columns([1, 1])
# with col1:
#     if st.button("Submit"):
#         if user_input:
#             try:
#                 # Pass the user input through the conversation chain
#                 response = st.session_state.conversational_chain(
#                     {"question": user_input, "chat_history": st.session_state.chat_history}
#                 )

#                 # Display response
#                 st.session_state.chat_history.append({"question": user_input, "answer": response['answer']})
#                 st.write(f"**Bot:** {response['answer']}")
#             except Exception as e:
#                 st.error(f"Error generating response: {e}")
#         else:
#             st.warning("Please enter a question!")

# with col2:
#     if st.button("Clear Chat"):
#         st.session_state.chat_history = []
#         st.success("Chat history cleared!")

# # if st.session_state.chat_history:
# #     for message in st.session_state.chat_history:
# #         st.markdown(user_template.format(msg=message['user']), unsafe_allow_html=True)
# #         st.markdown(bot_template.format(msg=message['bot']), unsafe_allow_html=True)

#         if message.get("context_docs"):
#             with st.expander("Source Documents"):
#                 for doc in message["context_docs"]:
#                     st.write(f"Source: {doc.metadata.get('source', 'N/A')}")
#                     st.write(doc.page_content)

#================================code with Next quetion button==============
# import streamlit as st
# from langchain.document_loaders import DirectoryLoader, PyPDFLoader
# from langchain.text_splitter import RecursiveCharacterTextSplitter
# from langchain_community.llms import Ollama
# from langchain.vectorstores import FAISS
# from langchain.prompts import ChatPromptTemplate, MessagesPlaceholder
# from langchain_core.runnables.history import RunnableWithMessageHistory
# from langchain_community.chat_message_histories import ChatMessageHistory
# from langchain.chains import create_retrieval_chain
# from langchain.chains.combine_documents import create_stuff_documents_chain
# from sentence_transformers import SentenceTransformer, util
# from langchain_core.chat_history import BaseChatMessageHistory
# from langchain.chains import create_history_aware_retriever
# from langchain.embeddings import HuggingFaceEmbeddings
# import os

# # Define template for bot and user messages
# bot_template = '''
# <div style="display: flex; align-items: center; margin-bottom: 10px;">
#     <div style="flex-shrink: 0; margin-right: 10px;">
#         <img src="https://uxwing.com/wp-content/themes/uxwing/download/communication-chat-call/answer-icon.png" 
#              style="max-height: 50px; max-width: 50px; border-radius: 50%; object-fit: cover;">
#     </div>
#     <div style="background-color: #f1f1f1; padding: 10px; border-radius: 10px; max-width: 75%; word-wrap: break-word; overflow-wrap: break-word;">
#         {msg}
#     </div>
# </div>
# '''

# user_template = '''
# <div style="display: flex; align-items: center; margin-bottom: 10px; justify-content: flex-end;">
#     <div style="flex-shrink: 0; margin-left: 10px;">
#         <img src="https://cdn.iconscout.com/icon/free/png-512/free-q-characters-character-alphabet-letter-36051.png?f=webp&w=512" 
#              style="max-height: 50px; max-width: 50px; border-radius: 50%; object-fit: cover;">
#     </div>    
#     <div style="background-color: #007bff; color: white; padding: 10px; border-radius: 10px; max-width: 75%; word-wrap: break-word; overflow-wrap: break-word;">
#         {msg}
#     </div>
# </div>
# '''

# # Function to prepare and split documents from a specified directory
# def prepare_and_split_docs(pdf_directory):
#     loader = DirectoryLoader(pdf_directory, glob="**/*.pdf", loader_cls=PyPDFLoader)
#     documents = loader.load()
#     splitter = RecursiveCharacterTextSplitter.from_tiktoken_encoder(
#         chunk_size=512,
#         chunk_overlap=256,
#         disallowed_special=(),
#         separators=["\n\n", "\n", " "]
#     )
#     return splitter.split_documents(documents)

# # Function to ingest documents into the vector database
# def ingest_into_vectordb(split_docs):
#     embeddings = HuggingFaceEmbeddings(model_name='sentence-transformers/all-MiniLM-L6-v2')
#     return FAISS.from_documents(split_docs, embeddings)

# def get_conversation_chain(retriever):
#     llm = Ollama(model="meta-llama/Llama-3.2-3B")
#     contextualize_q_system_prompt = (
#         "Given the chat history and the latest user question, "
#         "provide a response that directly addresses the user's query based on the provided documents. "
#         "If no relevant answer is found, respond with: "
#         "'I'm sorry, but I couldn’t find an answer to your question. Could you rephrase or provide more details?' "
#         "Do not rephrase the question or ask follow-up questions."
#     )
#     contextualize_q_prompt = ChatPromptTemplate.from_messages(
#         [
#             ("system", contextualize_q_system_prompt),
#             MessagesPlaceholder("chat_history"),
#             ("human", "{input}"),
#         ]
#     )
#     history_aware_retriever = create_history_aware_retriever(
#         llm, retriever, contextualize_q_prompt
#     )
#     system_prompt = (
#         "As a chat assistant, provide accurate and relevant information based on the provided document in 2-3 sentences. "
#         "Answer should be limited to 50 words and 2-3 sentences. If no relevant information is found, respond with: "
#         "'I'm sorry, but I couldn’t find an answer to your question. Could you rephrase or provide more details?' "
#         "{context}"
#     )
#     qa_prompt = ChatPromptTemplate.from_messages(
#         [
#             ("system", system_prompt),
#             MessagesPlaceholder("chat_history"),
#             ("human", "{input}"),
#         ]
#     )
#     question_answer_chain = create_stuff_documents_chain(llm, qa_prompt)
#     rag_chain = create_retrieval_chain(history_aware_retriever, question_answer_chain)

#     store = {}

#     def get_session_history(session_id: str) -> BaseChatMessageHistory:
#         if session_id not in store:
#             store[session_id] = ChatMessageHistory()
#         return store[session_id]

#     return RunnableWithMessageHistory(
#         rag_chain,
#         get_session_history,
#         input_messages_key="input",
#         history_messages_key="chat_history",
#         output_messages_key="answer",
#     )

# # Main Streamlit app
# logopath = r"C:\Users\Akshada.P\OneDrive - SAKSOFT LIMITED\Documents\SAKSOFT Internal BOT\saksoftimg.png"
# st.image(logopath, width=200)
# st.title("AssureBot")
# st.write("Welcome! How can I help you today?")

# # Define the path to your PDF directory
# pdf_directory = "Data/SalesDocs"  # Update with your actual PDF folder path

# # Prepare and ingest documents into the vector store
# split_docs = prepare_and_split_docs(pdf_directory)
# vector_db = ingest_into_vectordb(split_docs)
# retriever = vector_db.as_retriever()

# # Initialize the conversation chain
# if 'conversational_chain' not in st.session_state:
#     st.session_state.conversational_chain = get_conversation_chain(retriever)

# if 'chat_history' not in st.session_state:
#     st.session_state.chat_history = []

# if 'user_input' not in st.session_state:
#     st.session_state.user_input = ""

# # Chat input
# user_input = st.text_input("Ask a question about the documents:", value=st.session_state.user_input)

# # Buttons for submitting and refreshing
# # Buttons for submitting and refreshing
# col1, col2 = st.columns([1, 1])
# with col1:
#     if st.button("Submit"):
#         user_input_lower = user_input.lower()  # Normalize input for checking

#         # Store the user input in session state
#         st.session_state.user_input = user_input

#         # Respond to greetings
#         if any(greet in user_input_lower for greet in ["hi", "hello", "hey"]):
#             response = "Hello! How can I assist you?"
#             st.session_state.chat_history.append({"user": user_input, "bot": response, "context_docs": []})

#         # Respond to thanks
#         elif any(thank in user_input_lower for thank in ["thank you", "thanks", "thank you very much", "thx"]):
#             response = "Thank you! Let me know if you have any more queries."
#             st.session_state.chat_history.append({"user": user_input, "bot": response, "context_docs": []})

#         # Respond to yes/no
#         elif any(yn in user_input_lower for yn in ["yes", "no"]):
#             response = "I understand. If you have any other questions, feel free to ask!"
#             st.session_state.chat_history.append({"user": user_input, "bot": response, "context_docs": []})

#         # Process document queries
#         elif user_input and 'conversational_chain' in st.session_state:
#             session_id = "abc123"  # Static session ID for this demo; you can make it dynamic if needed
#             response = st.session_state.conversational_chain.invoke({"input": user_input}, config={"configurable": {"session_id": session_id}})
#             context_docs = response.get('context', [])

#             # Check if the response is the default message
#             if response['answer'] == "I'm sorry, but I couldn’t find an answer to your question. Could you rephrase or provide more details?":
#                 st.session_state.chat_history.append({"user": user_input, "bot": response['answer'], "context_docs": []})
#             else:
#                 st.session_state.chat_history.append({"user": user_input, "bot": response['answer'], "context_docs": context_docs})

# with col2:
#     if st.button("Next Question"):
#         st.session_state.user_input = ""  # Clear the input without refreshing the app

# # Display chat history
# if st.session_state.chat_history:
#     for message in st.session_state.chat_history:
#         st.markdown(user_template.format(msg=message['user']), unsafe_allow_html=True)
#         st.markdown(bot_template.format(msg=message['bot']), unsafe_allow_html=True)

#         # Display the source documents if available
#         if message.get('context_docs'):
#             with st.expander("Source Documents"):
#                 for doc in message['context_docs']:
#                     st.write(f"Source: {doc.metadata['source']}")
#                     st.write(doc.page_content)
# import streamlit as st
# from langchain.document_loaders import DirectoryLoader, PyPDFLoader
# from langchain.text_splitter import RecursiveCharacterTextSplitter
# from langchain_community.llms import Ollama
# from langchain.vectorstores import FAISS
# from langchain.prompts import ChatPromptTemplate, MessagesPlaceholder
# from langchain_core.runnables.history import RunnableWithMessageHistory
# from langchain_community.chat_message_histories import ChatMessageHistory
# from langchain.chains import create_retrieval_chain
# from langchain.chains.combine_documents import create_stuff_documents_chain
# from sentence_transformers import SentenceTransformer, util
# from langchain_core.chat_history import BaseChatMessageHistory
# from langchain.chains import create_history_aware_retriever
# from langchain.embeddings import HuggingFaceEmbeddings
# import os

# # Define template for bot and user messages
# bot_template = '''
# <div style="display: flex; align-items: center; margin-bottom: 10px;">
#     <div style="flex-shrink: 0; margin-right: 10px;">
#         <img src="https://uxwing.com/wp-content/themes/uxwing/download/communication-chat-call/answer-icon.png" 
#              style="max-height: 50px; max-width: 50px; border-radius: 50%; object-fit: cover;">
#     </div>
#     <div style="background-color: #f1f1f1; padding: 10px; border-radius: 10px; max-width: 75%; word-wrap: break-word; overflow-wrap: break-word;">
#         {msg}
#     </div>
# </div>
# '''

# user_template = '''
# <div style="display: flex; align-items: center; margin-bottom: 10px; justify-content: flex-end;">
#     <div style="flex-shrink: 0; margin-left: 10px;">
#         <img src="https://cdn.iconscout.com/icon/free/png-512/free-q-characters-character-alphabet-letter-36051.png?f=webp&w=512" 
#              style="max-height: 50px; max-width: 50px; border-radius: 50%; object-fit: cover;">
#     </div>    
#     <div style="background-color: #007bff; color: white; padding: 10px; border-radius: 10px; max-width: 75%; word-wrap: break-word; overflow-wrap: break-word;">
#         {msg}
#     </div>
# </div>
# '''

# # Function to prepare and split documents from a specified directory
# def prepare_and_split_docs(pdf_directory):
#     loader = DirectoryLoader(pdf_directory, glob="**/*.pdf", loader_cls=PyPDFLoader)
#     documents = loader.load()
#     splitter = RecursiveCharacterTextSplitter.from_tiktoken_encoder(
#         chunk_size=512,
#         chunk_overlap=256,
#         disallowed_special=(),
#         separators=["\n\n", "\n", " "]
#     )
#     return splitter.split_documents(documents)

# # Function to ingest documents into the vector database
# def ingest_into_vectordb(split_docs):
#     embeddings = HuggingFaceEmbeddings(model_name='sentence-transformers/all-MiniLM-L6-v2')
#     return FAISS.from_documents(split_docs, embeddings)

# def get_conversation_chain(retriever):
#     llm = Ollama(model="llama3.2")
#     contextualize_q_system_prompt = (
#         "Given the chat history and the latest user question, "
#         "provide a response that directly addresses the user's query based on the provided documents. "
#         "If no relevant answer is found, respond with: "
#         "'I'm sorry, but I couldn’t find an answer to your question. Could you rephrase or provide more details?' "
#         "Do not rephrase the question or ask follow-up questions."
#     )
#     contextualize_q_prompt = ChatPromptTemplate.from_messages(
#         [
#             ("system", contextualize_q_system_prompt),
#             MessagesPlaceholder("chat_history"),
#             ("human", "{input}"),
#         ]
#     )
#     history_aware_retriever = create_history_aware_retriever(
#         llm, retriever, contextualize_q_prompt
#     )
#     system_prompt = (
#         "As a chat assistant, provide accurate and relevant information based on the provided document in 2-3 sentences. "
#         "Answer should be contextualize and relevant words and  sentences. If no relevant information is found, respond with: "
#         "'I'm sorry, but I couldn’t find an answer to your question. Could you rephrase or provide more details?' "
#         "{context}"
#     )
#     qa_prompt = ChatPromptTemplate.from_messages(
#         [
#             ("system", system_prompt),
#             MessagesPlaceholder("chat_history"),
#             ("human", "{input}"),
#         ]
#     )
#     question_answer_chain = create_stuff_documents_chain(llm, qa_prompt)
#     rag_chain = create_retrieval_chain(history_aware_retriever, question_answer_chain)

#     store = {}

#     def get_session_history(session_id: str) -> BaseChatMessageHistory:
#         if session_id not in store:
#             store[session_id] = ChatMessageHistory()
#         return store[session_id]

#     return RunnableWithMessageHistory(
#         rag_chain,
#         get_session_history,
#         input_messages_key="input",
#         history_messages_key="chat_history",
#         output_messages_key="answer",
#     )

# # Main Streamlit app
# st.set_page_config(page_title="AssureBot", layout="wide")
# logopath = r"C:\Users\Akshada.P\OneDrive - SAKSOFT LIMITED\Documents\SAKSOFT Internal BOT\saksoftimg.png"
# st.image(logopath, width=200)
# st.title("AssureBot")
# st.write("Welcome! How can I help you today?")

# # Define the path to your PDF directory
# pdf_directory = "Data/SalesDocs"  # Update with your actual PDF folder path

# # Prepare and ingest documents into the vector store
# split_docs = prepare_and_split_docs(pdf_directory)
# vector_db = ingest_into_vectordb(split_docs)
# retriever = vector_db.as_retriever()

# # Initialize the conversation chain
# if 'conversational_chain' not in st.session_state:
#     st.session_state.conversational_chain = get_conversation_chain(retriever)

# if 'chat_history' not in st.session_state:
#     st.session_state.chat_history = []

# if 'user_input' not in st.session_state:
#     st.session_state.user_input = ""

# # Chat input
# user_input = st.text_input("Ask a question about the documents:", value=st.session_state.user_input)

# # Buttons for submitting and next question
# col1, col2 = st.columns([1, 1])
# with col1:
#     if st.button("Submit"):
#         user_input_lower = user_input.lower()  # Normalize input for checking

#         # Store the user input in session state
#         st.session_state.user_input = user_input

#         # Respond to greetings
#         if any(greet in user_input_lower for greet in ["hi", "hello", "hey"]):
#             response = "Hello! How can I assist you?"
#             st.session_state.chat_history.append({"user": user_input, "bot": response, "context_docs": []})

#         # Respond to thanks
#         elif any(thank in user_input_lower for thank in ["thank you", "thanks", "thank you very much", "thx"]):
#             response = "Thank you! Let me know if you have any more queries."
#             st.session_state.chat_history.append({"user": user_input, "bot": response, "context_docs": []})

#         # Respond to yes/no
#         elif any(yn in user_input_lower for yn in ["yes", "no"]):
#             response = "I understand. If you have any other questions, feel free to ask!"
#             st.session_state.chat_history.append({"user": user_input, "bot": response, "context_docs": []})

#         # Process document queries
#         elif user_input and 'conversational_chain' in st.session_state:
#             session_id = "abc123"  # Static session ID for this demo; you can make it dynamic if needed
#             response = st.session_state.conversational_chain.invoke({"input": user_input}, config={"configurable": {"session_id": session_id}})
#             context_docs = response.get('context', [])

#             # Check if the response is the default message
#             if response['answer'] == "I'm sorry, but I couldn’t find an answer to your question. Could you rephrase or provide more details?":
#                 st.session_state.chat_history.append({"user": user_input, "bot": response['answer'], "context_docs": []})
#             else:
#                 st.session_state.chat_history.append({"user": user_input, "bot": response['answer'], "context_docs": context_docs})

# with col2:
#     if st.button("Next Question"):
#         st.session_state.user_input = ""  # Clear the input

# # Display chat history
# if st.session_state.chat_history:
#     for message in st.session_state.chat_history:
#         st.markdown(user_template.format(msg=message['user']), unsafe_allow_html=True)
#         st.markdown(bot_template.format(msg=message['bot']), unsafe_allow_html=True)

#         # Display the source documents if available
#         if message.get('context_docs'):
#             with st.expander("Source Documents"):
#                 for doc in message['context_docs']:
#                     st.write(f"Source: {doc.metadata['source']}")
#                     st.write(doc.page_content)
# import streamlit as st
# from langchain.document_loaders import DirectoryLoader, PyPDFLoader
# from langchain.text_splitter import RecursiveCharacterTextSplitter
# from langchain_community.llms import Ollama
# from langchain.vectorstores import FAISS
# from langchain.prompts import ChatPromptTemplate, MessagesPlaceholder
# from langchain_core.runnables.history import RunnableWithMessageHistory
# from langchain_community.chat_message_histories import ChatMessageHistory
# from langchain.chains import create_retrieval_chain
# from langchain.chains.combine_documents import create_stuff_documents_chain
# from sentence_transformers import SentenceTransformer, util
# from langchain_core.chat_history import BaseChatMessageHistory
# from langchain.chains import create_history_aware_retriever
# from langchain.embeddings import HuggingFaceEmbeddings


# # Templates for user and bot messages
# bot_template = '''
# <div style="display: flex; align-items: center; margin-bottom: 10px;">
#     <div style="flex-shrink: 0; margin-right: 10px;">
#         <img src="https://uxwing.com/wp-content/themes/uxwing/download/communication-chat-call/answer-icon.png" 
#              style="max-height: 50px; max-width: 50px; border-radius: 50%; object-fit: cover;">
#     </div>
#     <div style="background-color: #f1f1f1; padding: 10px; border-radius: 10px; max-width: 75%; word-wrap: break-word; overflow-wrap: break-word;">
#         {msg}
#     </div>
# </div>
# '''

# user_template = '''
# <div style="display: flex; align-items: center; margin-bottom: 10px; justify-content: flex-end;">
#     <div style="flex-shrink: 0; margin-left: 10px;">
#         <img src="https://cdn.iconscout.com/icon/free/png-512/free-q-characters-character-alphabet-letter-36051.png?f=webp&w=512" 
#              style="max-height: 50px; max-width: 50px; border-radius: 50%; object-fit: cover;">
#     </div>    
#     <div style="background-color: #007bff; color: white; padding: 10px; border-radius: 10px; max-width: 75%; word-wrap: break-word; overflow-wrap: break-word;">
#         {msg}
#     </div>
# </div>
# '''

# # Function to prepare and split documents
# def prepare_and_split_docs(pdf_directory):
#     loader = DirectoryLoader(pdf_directory, glob="**/*.pdf", loader_cls=PyPDFLoader)
#     documents = loader.load()
#     splitter = RecursiveCharacterTextSplitter.from_tiktoken_encoder(
#         chunk_size=512,
#         chunk_overlap=256,
#         disallowed_special=(),
#         separators=["\n\n", "\n", " "]
#     )
#     split_docs = splitter.split_documents(documents)
#     return split_docs

# # Function to ingest documents into vector store
# def ingest_into_vectordb(split_docs):
#     embeddings = HuggingFaceEmbeddings(model_name='sentence-transformers/all-MiniLM-L6-v2')
#     db = FAISS.from_documents(split_docs, embeddings)
#     return db

# # Function to create a conversation chain
# def get_conversation_chain(retriever):
#     llm = Ollama(model="llama3.2")
#     contextualize_q_system_prompt = (
#         "Given the chat history and the latest user question, "
#         "provide a response based on the documents. If no answer is found, "
#         "respond: 'I'm sorry, but I couldn’t find an answer. Could you rephrase or provide more details?'"
#     )
#     contextualize_q_prompt = ChatPromptTemplate.from_messages(
#         [
#             ("system", contextualize_q_system_prompt),
#             MessagesPlaceholder("chat_history"),
#             ("human", "{input}"),
#         ]
#     )
#     history_aware_retriever = create_history_aware_retriever(
#         llm, retriever, contextualize_q_prompt
#     )
#     system_prompt = (
#         "As a chat assistant, provide accurate and relevant information based on the provided document in 2-3 sentences. "
#         "Answer should be correct to the point short and brief for given question. If no relevant information is found, respond with: "
#         "'I'm sorry, but I couldn’t find an answer to your question. Could you rephrase or provide more details?' "
#         "{context}"
#     )
#     qa_prompt = ChatPromptTemplate.from_messages(
#         [
#             ("system", system_prompt),
#             MessagesPlaceholder("chat_history"),
#             ("human", "{input}"),
#         ]
#     )
#     question_answer_chain = create_stuff_documents_chain(llm, qa_prompt)
#     rag_chain = create_retrieval_chain(history_aware_retriever, question_answer_chain)

#     store = {}

#     def get_session_history(session_id: str):
#         if session_id not in store:
#             store[session_id] = ChatMessageHistory()
#         return store[session_id]

#     conversational_rag_chain = RunnableWithMessageHistory(
#         rag_chain,
#         get_session_history,
#         input_messages_key="input",
#         history_messages_key="chat_history",
#         output_messages_key="answer",
#     )
#     return conversational_rag_chain

# # Main Streamlit app
# logopath = r"C:\Users\Akshada.P\OneDrive - SAKSOFT LIMITED\Documents\SAKSOFT Internal BOT\saksoftimg.png"
# st.image(logopath, width=200)
# st.title("Q&A Bot")
# st.write("Welcome! How can I help you today?")

# pdf_directory = "Data/SalesDocs"  # Update to your actual folder path

# # Prepare documents and ingest into vector store
# split_docs = prepare_and_split_docs(pdf_directory)
# vector_db = ingest_into_vectordb(split_docs)
# retriever = vector_db.as_retriever()

# # Initialize conversation chain
# conversational_chain = get_conversation_chain(retriever)
# st.session_state.conversational_chain = conversational_chain

# if 'chat_history' not in st.session_state:
#     st.session_state.chat_history = []

# # Chat input
# user_input = st.text_input("Ask a question about the documents:")
# # Buttons for submitting and refreshing
# col1, col2 = st.columns([1, 1])
# with col1:
#     if st.button("Submit"):
#         user_input_lower = user_input.lower()  # Normalize input for checking

#         # Store the user input in session state
#         st.session_state.user_input = user_input

#         # Respond to greetings
#         if any(greet in user_input_lower for greet in ["hi", "hello", "hey"]):
#             response = "Hello! How can I assist you?"
#             st.session_state.chat_history.append({"user": user_input, "bot": response, "context_docs": []})

#         # Respond to thanks
#         elif any(thank in user_input_lower for thank in ["thank you", "thanks", "thank you very much", "thx"]):
#             response = "Thank you! Let me know if you have any more queries."
#             st.session_state.chat_history.append({"user": user_input, "bot": response, "context_docs": []})

#         # Respond to yes/no
#         elif any(yn in user_input_lower for yn in ["yes", "no"]):
#             response = "I understand. If you have any other questions, feel free to ask!"
#             st.session_state.chat_history.append({"user": user_input, "bot": response, "context_docs": []})

#         # Process document queries
#         elif user_input and 'conversational_chain' in st.session_state:
#             session_id = "abc123"  # Static session ID for this demo; you can make it dynamic if needed
#             response = st.session_state.conversational_chain.invoke({"input": user_input}, config={"configurable": {"session_id": session_id}})
#             context_docs = response.get('context', [])

#             # Check if the response is the default message
#             if response['answer'] == "I'm sorry, but I couldn’t find an answer to your question. Could you rephrase or provide more details?":
#                 st.session_state.chat_history.append({"user": user_input, "bot": response['answer'], "context_docs": []})
#             else:
#                 st.session_state.chat_history.append({"user": user_input, "bot": response['answer'], "context_docs": context_docs})

# with col2:
#     if st.button("Refresh"):
#         st.session_state.user_input = ""  # Clear the input
#         st.session_state.chat_history = []  # Clear chat history

# # Display chat history
# for message in st.session_state.chat_history:
#     st.markdown(user_template.format(msg=message['user']), unsafe_allow_html=True)
#     st.markdown(bot_template.format(msg=message['bot']), unsafe_allow_html=True)

#     if message.get("context_docs"):
#         with st.expander("Source Documents"):
#             for doc in message["context_docs"]:
#                 st.write(f"Source: {doc.metadata.get('source', 'N/A')}")
#                 st.write(doc.page_content)